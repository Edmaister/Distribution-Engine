from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


TEAL = RGBColor(0x00, 0x89, 0x7B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x72, 0x80)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0


def add_two_column(prs: Presentation, title: str, left: list[str], right: list[str], left_title: str, right_title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.9), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(5.9), Inches(5.5))

    for box, heading, lines in ((left_box, left_title, left), (right_box, right_title, right)):
        tf = box.text_frame
        tf.clear()
        h = tf.paragraphs[0]
        h.text = heading
        h.font.bold = True
        h.font.size = Pt(20)
        h.font.color.rgb = TEAL
        for line in lines:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = DARK


def style(prs: Presentation) -> None:
    for slide in prs.slides:
        if slide.shapes.title is not None:
            t = slide.shapes.title.text_frame
            for p in t.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(34)
                    r.font.bold = True
                    r.font.color.rgb = TEAL


def build() -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        "Referral Platform Executive Architecture",
        "Capability, Data Structure, Integration and Operating Model\nPrepared for executive stakeholders",
    )

    add_bullets(
        prs,
        "Executive Summary",
        [
            "The platform is a persona-safe referral engine combining referral tracking, rewards, missions, badges, and leaderboard services.",
            "It exposes consolidated dashboard APIs for referrer and referee journeys, with privacy controls at API contract level.",
            "The architecture is cloud-ready: FastAPI services, PostgreSQL data layer, Kafka event integration, and Prometheus observability.",
            "The frontend integration uses a normalized adapter layer and compliance-safe UX copy to reduce data leakage risk and support regulated communication.",
        ],
    )

    add_bullets(
        prs,
        "Business Capabilities Delivered",
        [
            "Referral lifecycle: bootstrap, code issuance, acceptance, progress tracking, completion.",
            "Reward management: totals, pending, potential, next eligible, line-item explanations.",
            "Mission framework: core, boost, and milestone mission categories with progress and payout semantics.",
            "Engagement layers: badges, leaderboard rank/tier progression, recommendations.",
            "Governance: disclosures, compliance metadata, role-based visibility and consent handling.",
        ],
    )

    add_two_column(
        prs,
        "Reference Architecture",
        left_title="Application & API Layer",
        right_title="Data, Integration & Ops Layer",
        left=[
            "FastAPI service with domain routers: dashboard, missions, rewards, referrals, badges, leaderboards.",
            "Adapter boundary (`ReferralApiAdapter`) keeps clients stable while backend evolves.",
            "Persona-safe endpoints: `/v1/referrers/{ucn}/dashboard`, `/v1/referrals/{id}/referrer-view`, `/v1/dashboard/referee/{id}`.",
            "Jinja/Front-end integration consumes normalized contracts and compliance copy catalog.",
        ],
        right=[
            "PostgreSQL persistence with migrations and seed data for mission definitions and scoring rules.",
            "Kafka integration for event-driven processing and replay/worker operations.",
            "Health model: `/healthz` for liveness, `/readyz` for dependency-aware readiness.",
            "Prometheus metrics: request count, latency histogram, DB/Kafka readiness gauges.",
        ],
    )

    add_bullets(
        prs,
        "Core API Contract (Referrer)",
        [
            "`GET /v1/referrers/{referrer_ucn}/dashboard`: summary totals, rewards, grouped missions, badges, leaderboard, referral cards.",
            "`GET /v1/missions/referrer/{referrer_ucn}`: grouped mission payload (`core`, `boost`, `milestone`) and mission metadata.",
            "`GET /v1/rewards/summary/referrers/{referrer_ucn}`: earned/pending/potential totals, counts, disclosures, compliance, optional line items.",
            "`GET /v1/referrals/{referral_track_id}/referrer-view`: persona-safe progress and reward state for detail screens.",
        ],
    )

    add_bullets(
        prs,
        "Core API Contract (Referee)",
        [
            "`GET /v1/dashboard/referee/{referral_track_id}`: journey milestones, conditions and potential reward.",
            "`GET /v1/referrals/{referral_track_id}/referee-view`: detailed journey state.",
            "`GET /v1/rewards/referee/{referral_track_id}`: referee reward summary/list.",
            "`GET /v1/rewards/{reward_id}/explanation`: reason, trigger milestone, policy and outstanding conditions.",
        ],
    )

    add_two_column(
        prs,
        "Canonical Data Structure",
        left_title="Domain Entities",
        right_title="Key Relationships",
        left=[
            "Referrer profile and referral code",
            "Referral instance (track ID, product, status, progress, milestones)",
            "Reward summary (totals + line items + compliance metadata)",
            "Mission item (category, progress, reward amount, status)",
            "Badge and leaderboard entries (score and rank metadata)",
        ],
        right=[
            "One referrer has many referral instances.",
            "Each referral instance maps to reward summary slices and mission progress.",
            "Missions grouped by `core/boost/milestone` with payout semantics.",
            "Leaderboard aggregates referral activity into rank tier and next-rank deltas.",
            "Compliance/disclosure metadata attached to reward and mission responses.",
        ],
    )

    add_bullets(
        prs,
        "Data & Privacy Controls",
        [
            "Persona visibility is enforced in API responses, not inferred in templates.",
            "Referrer views exclude sensitive referee action hints; referee views include actionable journey detail.",
            "Consent schema supports auditability (`acceptedTerms`, timestamps, permission confirmation).",
            "Frontend normalizes data and applies policy-safe copy to avoid exposing disallowed detail.",
        ],
    )

    add_bullets(
        prs,
        "End-to-End Integration Flow",
        [
            "1) User context resolved in frontend; referrals bootstrap/code endpoints establish campaign participation.",
            "2) Frontend calls dashboard/reward/mission APIs through adapter and normalization services.",
            "3) UI renders consolidated journey (overview, progress, boost rewards, detail, leaderboard).",
            "4) Reward explanations and mission progress reinforce transparency and compliance-ready communication.",
        ],
    )

    add_bullets(
        prs,
        "Operational Architecture",
        [
            "Resilience: dependency health caching and fail-fast DB/Kafka probes reduce readiness-induced outages.",
            "Observability: Prometheus metrics for request volume, latency, and dependency availability.",
            "Deployment posture: container-friendly FastAPI services suited for EB/ECS/Kubernetes rollout patterns.",
            "Static asset versioning and UI contract docs support safer frontend release management.",
        ],
    )

    add_bullets(
        prs,
        "Current Strengths and Improvement Levers",
        [
            "Strengths: clear domain separation, explicit API contracts, robust compliance-aware UX integration.",
            "Levers: standardize schema governance, introduce event replay dashboards, and formalize API versioning policy.",
            "Levers: add architecture decision records and data lineage mapping for audit and change control.",
            "Levers: deepen SLO/SLA reporting and business KPI telemetry (conversion, payout efficiency, mission completion).",
        ],
    )

    add_bullets(
        prs,
        "Executive Roadmap (90-Day View)",
        [
            "Phase 1: Stabilize contract governance and publish platform capability scorecard.",
            "Phase 2: Expand analytics layer (funnel, cohort, mission efficacy, rank progression).",
            "Phase 3: Introduce partner-ready integration package (API docs, event schema pack, controls checklist).",
            "Phase 4: Scale readiness (performance testing, resilience drills, compliance evidence automation).",
        ],
    )

    add_title_slide(
        prs,
        "Appendix / Q&A",
        "Detailed API contract matrix, schema mappings, and integration traceability available on request.",
    )

    style(prs)
    out_path = r"c:\Projects\Referral Engine\Referral_Engine_Executive_Architecture.pptx"
    prs.save(out_path)
    print(out_path)


if __name__ == "__main__":
    build()

